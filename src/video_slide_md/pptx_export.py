# FILE: src/video_slide_md/pptx_export.py
# VERSION: 0.1.0
# START_MODULE_CONTRACT
#   PURPOSE: Export slides.json to PPTX presentation with slide images and speaker notes
#   SCOPE: Generate .pptx with full-slide screenshots, formatted transcript as speaker notes,
#          subtitle cues grouped into readable author paragraphs
#   DEPENDS: models, pathlib, loguru, python-pptx
#   LINKS: M-PPTX-EXPORT
#   ROLE: RUNTIME
#   MAP_MODE: EXPORTS
# END_MODULE_CONTRACT
#
# START_MODULE_MAP
#   export_to_pptx - write .pptx with all slides (screenshots + speaker notes)
#   export_slides_to_pptx - convenience: read slides.json, write .pptx
# END_MODULE_MAP

from __future__ import annotations

from pathlib import Path

from loguru import logger

from video_slide_md.models import SlideSegment, SlidesDocument, SubtitleCue


def export_to_pptx(
    document: SlidesDocument,
    output_path: Path,
    slides_dir: str | Path = "slides",
    title: str = "Presentation",
    notes_mode: str = "basic",
) -> Path:
    # START_CONTRACT: export_to_pptx
    #   PURPOSE: Write .pptx with slide images as full-slide pictures,
    #            speaker notes containing processed transcript
    #   INPUTS: {
    #       document: SlidesDocument — slides with images + cues,
    #       output_path: Path — target .pptx path,
    #       slides_dir: str | Path — relative dir for slide images,
    #       title: str — presentation title,
    #       notes_mode: str — "basic" or "llm"
    #   }
    #   OUTPUTS: Path to .pptx file
    #   SIDE_EFFECTS: Creates .pptx file
    #   LINKS: M-PPTX-EXPORT
    # END_CONTRACT: export_to_pptx

    # START_BLOCK_PPTX_BUILD
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(12192000)    # 16:9 widescreen
    prs.slide_height = Emu(6858000)

    slides_root = Path(slides_dir)

    for seg in document.slides:
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # START_BLOCK_ADD_IMAGE
        img_path = slides_root / seg.image if seg.image else None
        if img_path and img_path.is_file():
            slide.shapes.add_picture(
                str(img_path),
                Emu(0), Emu(0),
                prs.slide_width, prs.slide_height,
            )
        # END_BLOCK_ADD_IMAGE

        # START_BLOCK_ADD_NOTES
        notes_text = _format_slide_notes(seg, notes_mode=notes_mode)
        if notes_text.strip():
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text.strip()
        # END_BLOCK_ADD_NOTES

    prs.save(str(output_path))
    logger.info(
        f"[PptxExport][export_to_pptx] PPTX written | "
        f"path={output_path} slides={len(document.slides)}"
    )
    return output_path
    # END_BLOCK_PPTX_BUILD


def _format_slide_notes(seg: SlideSegment, notes_mode: str = "basic") -> str:
    # START_CONTRACT: _format_slide_notes
    #   PURPOSE: Format slide transcript into speaker notes via notes_processor
    #   INPUTS: { seg: SlideSegment, notes_mode: str }
    #   OUTPUTS: str — formatted notes text
    #   SIDE_EFFECTS: none (basic mode); may call LLM in llm mode
    #   LINKS: M-PPTX-EXPORT
    # END_CONTRACT: _format_slide_notes

    from video_slide_md.notes_processor import process_notes

    lines: list[str] = []

    # Timestamp header
    lines.append(f"[ {_fmt_time(seg.start)} – {_fmt_time(seg.end)} ]")
    lines.append("")

    # Process notes
    notes = process_notes(seg, mode=notes_mode)
    if notes:
        lines.append(notes)

    return "\n".join(lines).strip()


def _group_cues(
    cues: list[SubtitleCue],
    gap_threshold: float = 2.0,
) -> list[list[SubtitleCue]]:
    # START_CONTRACT: _group_cues
    #   PURPOSE: Group consecutive subtitle cues into paragraphs where gap <= threshold
    #   INPUTS: { cues: list[SubtitleCue], gap_threshold: float }
    #   OUTPUTS: list[list[SubtitleCue]]
    #   SIDE_EFFECTS: none
    #   LINKS: M-PPTX-EXPORT
    # END_CONTRACT: _group_cues

    if not cues:
        return []

    paragraphs: list[list[SubtitleCue]] = [[cues[0]]]
    for cue in cues[1:]:
        gap = cue.start - paragraphs[-1][-1].end
        if gap <= gap_threshold:
            paragraphs[-1].append(cue)
        else:
            paragraphs.append([cue])
    return paragraphs


def _fmt_time(seconds: float) -> str:
    h = int(seconds // 3600)
    m = int((seconds % 3600) // 60)
    s = int(seconds % 60)
    if h > 0:
        return f"{h}:{m:02d}:{s:02d}"
    return f"{m}:{s:02d}"
