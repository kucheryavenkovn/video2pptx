# FILE: tests/e2e/test_mcp_startup_characterization.py
# VERSION: 1.0.0
# START_MODULE_CONTRACT
#   PURPOSE: Characterize real GUI subprocess startup and MCP transport without mocks.
#   SCOPE: Process lifecycle, fresh port discovery, health, initialize, tools/list, initial UI state.
#   DEPENDS: pytest, tools/mcp_e2e_runner.py, PySide6
#   LINKS: V-M-REF-CHAR-TESTS, E2E-001, M-GUI-PROJECT-CTRL
#   ROLE: TEST
#   MAP_MODE: LOCALS
# END_MODULE_CONTRACT

from __future__ import annotations

import json
import re

from tools.mcp_e2e_runner import GuiProcessHarness, McpClient


def _call_and_wait(client, tool, arguments, timeout=30):
    queued = McpClient.result_data(client.tool_call(tool, arguments))
    assert queued["status"] == "queued"
    completed = client.wait_operation(queued["operation_id"], timeout=timeout)
    assert completed["status"] == "succeeded", completed
    return completed


def _normalized_timeline_state(timeline):
    tracks = timeline["tracks"]
    return {
        "slides": [
            {
                key: clip.get(key)
                for key in (
                    "uid",
                    "index",
                    "start_sec",
                    "end_sec",
                    "image_path",
                    "manual",
                    "transcript",
                )
            }
            for clip in tracks.get("slides", {}).get("clips", [])
        ],
        "subtitles": [
            {
                key: clip.get(key)
                for key in ("start_sec", "end_sec", "text")
            }
            for clip in tracks.get("subtitles", {}).get("clips", [])
        ],
        "scores": [
            {
                key: clip.get(key)
                for key in ("start_sec", "end_sec", "value", "method")
            }
            for clip in tracks.get("scores", {}).get("clips", [])
        ],
    }


def test_real_gui_mcp_startup(repo_dir, tmp_path):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        assert harness.process is not None
        assert harness.process.poll() is None
        assert harness.port_record is not None
        assert harness.port_record.pid == harness.process.pid
        assert harness.port_record.instance_id

        health = client.health()
        assert health["status"] == "ok"
        assert health["pid"] == harness.process.pid
        assert health["instance_id"] == harness.port_record.instance_id

        initialized = client.initialize()
        info = initialized["result"]["serverInfo"]
        assert info["name"] == "video2pptx"

        listed = client.tools_list()
        tool_names = {
            tool["name"] for tool in listed["result"]["tools"]
        }
        assert {
            "health",
            "get_ui_state",
            "get_project",
            "get_timeline",
            "project_create",
            "detect",
            "wait_operation",
        } <= tool_names

        ui = McpClient.result_data(client.tool_call("get_ui_state"))
        assert ui["window_title"].startswith("video2pptx")
        assert ui["busy"] is False
        assert ui["buttons"]["detect"]["enabled"] is False
        assert ui["buttons"]["auto_align"]["enabled"] is False
    finally:
        harness.stop()

    assert harness.process is not None
    assert harness.process.poll() is not None
    assert not harness.port_file.exists()


def test_real_gui_project_create_updates_model_and_window(repo_dir, tmp_path):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"
        queued = McpClient.result_data(
            client.tool_call(
                "project_create",
                {"path": str(project_parent), "name": "characterized"},
            )
        )
        assert queued["status"] == "queued"
        assert queued["operation_id"]

        completed = client.wait_operation(queued["operation_id"], timeout=10)
        assert completed["status"] == "succeeded"

        project = McpClient.result_data(client.tool_call("get_project"))
        assert project["name"] == "characterized"
        assert project["project_dir"] == str(project_parent / "characterized")
        assert project["slides_count"] == 0

        ui = McpClient.result_data(client.tool_call("get_ui_state"))
        assert "characterized" in ui["window_title"]
        assert ui["buttons"]["save"]["enabled"] is True
        assert (project_parent / "characterized" / "project.json").is_file()
    finally:
        harness.stop()


def test_real_gui_mcp_detect_updates_project_timeline_and_disk(
    repo_dir,
    tmp_path,
    synthetic_video_path,
    synthetic_subtitle_path,
):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"

        for tool, arguments in (
            (
                "project_create",
                {"path": str(project_parent), "name": "detect_characterized"},
            ),
            ("video_import", {"path": str(synthetic_video_path)}),
            ("subtitle_import", {"path": str(synthetic_subtitle_path)}),
        ):
            queued = McpClient.result_data(client.tool_call(tool, arguments))
            completed = client.wait_operation(queued["operation_id"], timeout=10)
            assert completed["status"] == "succeeded"

        queued = McpClient.result_data(
            client.tool_call("detect", {"confirm": True})
        )
        completed = client.wait_operation(queued["operation_id"], timeout=120)
        assert completed["status"] == "succeeded", completed

        project = McpClient.result_data(client.tool_call("get_project"))
        timeline = McpClient.result_data(client.tool_call("get_timeline"))
        project_dir = project_parent / "detect_characterized"

        assert project["pipeline_state"]["detect_done"] is True
        assert project["pipeline_state"]["align_done"] is False
        assert project["pipeline_state"]["notes_done"] is False
        assert project["slides_count"] > 0
        assert timeline["tracks"]["slides"]["clip_count"] == project["slides_count"]
        assert timeline["tracks"]["subtitles"]["clip_count"] == 4
        assert timeline["tracks"]["scores"]["clip_count"] > 0
        assert (project_dir / "project.json").is_file()
        assert (project_dir / "slides.json").is_file()
        assert not (project_dir / "deck.md").exists()
        assert not (project_dir / "deck.pptx").exists()

        ui = McpClient.result_data(client.tool_call("get_ui_state"))
        assert ui["buttons"]["auto_align"]["enabled"] is True
    finally:
        harness.stop()


def test_real_gui_mcp_quick_preview_is_idempotent_and_side_effect_free(
    repo_dir,
    tmp_path,
    synthetic_video_path,
):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"
        for tool, arguments in (
            (
                "project_create",
                {"path": str(project_parent), "name": "preview_characterized"},
            ),
            ("video_import", {"path": str(synthetic_video_path)}),
        ):
            queued = McpClient.result_data(client.tool_call(tool, arguments))
            completed = client.wait_operation(queued["operation_id"], timeout=10)
            assert completed["status"] == "succeeded"

        score_counts = []
        for _ in range(2):
            queued = McpClient.result_data(
                client.tool_call("quick_preview", {})
            )
            completed = client.wait_operation(queued["operation_id"], timeout=60)
            assert completed["status"] == "succeeded", completed
            timeline = McpClient.result_data(client.tool_call("get_timeline"))
            score_counts.append(timeline["tracks"]["scores"]["clip_count"])

        project = McpClient.result_data(client.tool_call("get_project"))
        project_dir = project_parent / "preview_characterized"
        assert project["pipeline_state"]["preview_done"] is True
        assert project["pipeline_state"]["detect_done"] is False
        assert project["slides_count"] == 0
        assert score_counts[0] > 0
        assert score_counts[1] == score_counts[0]
        assert not (project_dir / "slides.json").exists()
        assert not (project_dir / "slides").exists()
        assert not (project_dir / "deck.md").exists()
        assert not (project_dir / "deck.pptx").exists()
    finally:
        harness.stop()


def test_real_gui_mcp_auto_align_dry_run_apply_and_idempotency(
    repo_dir,
    tmp_path,
    synthetic_video_path,
    synthetic_subtitle_path,
):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"
        project_dir = project_parent / "align_characterized"
        _call_and_wait(
            client,
            "project_create",
            {"path": str(project_parent), "name": "align_characterized"},
        )
        _call_and_wait(
            client,
            "video_import",
            {"path": str(synthetic_video_path)},
        )
        _call_and_wait(
            client,
            "subtitle_import",
            {"path": str(synthetic_subtitle_path)},
        )
        _call_and_wait(client, "detect", {"confirm": True}, timeout=120)

        project_json = project_dir / "project.json"
        slides_json = project_dir / "slides.json"
        project_before = project_json.read_bytes()
        slides_before = slides_json.read_bytes()
        timeline_before = McpClient.result_data(
            client.tool_call("get_timeline")
        )

        dry_run = _call_and_wait(
            client,
            "auto_align",
            {"dry_run": True, "max_shift_sec": 3.0},
        )
        assert dry_run["result"]["dry_run"] is True
        assert project_json.read_bytes() == project_before
        assert slides_json.read_bytes() == slides_before
        assert not (project_dir / "alignment_report.json").exists()
        assert McpClient.result_data(
            client.tool_call("get_timeline")
        ) == timeline_before
        project = McpClient.result_data(client.tool_call("get_project"))
        assert project["pipeline_state"]["align_done"] is False

        _call_and_wait(
            client,
            "auto_align",
            {"dry_run": False, "max_shift_sec": 3.0, "confirm": True},
        )
        aligned_once = slides_json.read_bytes()
        document = json.loads(aligned_once)
        slides = document["slides"]
        assert slides
        assert [slide["index"] for slide in slides] == list(
            range(1, len(slides) + 1)
        )
        assert all(slide["start"] < slide["end"] for slide in slides)
        assert all(
            left["end"] <= right["start"]
            for left, right in zip(slides, slides[1:], strict=False)
        )
        assert (project_dir / "alignment_report.json").is_file()
        project = McpClient.result_data(client.tool_call("get_project"))
        assert project["pipeline_state"]["align_done"] is True
        assert project["slides_count"] == len(slides)

        _call_and_wait(
            client,
            "auto_align",
            {"dry_run": False, "max_shift_sec": 3.0, "confirm": True},
        )
        assert slides_json.read_bytes() == aligned_once
    finally:
        harness.stop()


def test_real_gui_mcp_slide_crud_keeps_four_views_consistent(
    repo_dir,
    tmp_path,
    synthetic_video_path,
):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"
        project_dir = project_parent / "crud_characterized"
        _call_and_wait(
            client,
            "project_create",
            {"path": str(project_parent), "name": "crud_characterized"},
        )
        _call_and_wait(
            client,
            "video_import",
            {"path": str(synthetic_video_path)},
        )
        _call_and_wait(client, "detect", {"confirm": True}, timeout=120)

        timeline = McpClient.result_data(client.tool_call("get_timeline"))
        original = timeline["tracks"]["slides"]["clips"]
        original_count = len(original)
        first = original[0]
        split_at = (first["start_sec"] + first["end_sec"]) / 2

        added = _call_and_wait(client, "slide_add", {"ts": split_at})
        uid = added["result"]["uid"]
        created = McpClient.result_data(
            client.tool_call("get_slide", {"uid": uid})
        )
        assert created["uid"] == uid
        assert created["manual"] is True

        project = McpClient.result_data(client.tool_call("get_project"))
        timeline = McpClient.result_data(client.tool_call("get_timeline"))
        document = json.loads((project_dir / "slides.json").read_text("utf-8"))
        project_doc = json.loads((project_dir / "project.json").read_text("utf-8"))
        assert project["slides_count"] == original_count + 1
        assert timeline["tracks"]["slides"]["clip_count"] == original_count + 1
        assert len(document["slides"]) == original_count + 1
        assert len(project_doc["slides"]) == original_count + 1

        resized_end = first["start_sec"] + (
            split_at - first["start_sec"]
        ) * 0.8
        _call_and_wait(
            client,
            "slide_resize",
            {"index": 1, "end": resized_end, "confirm": True},
        )
        _call_and_wait(
            client,
            "slide_move",
            {
                "uid": uid,
                "start": resized_end,
                "end": first["end_sec"],
                "confirm": True,
            },
        )
        moved = McpClient.result_data(
            client.tool_call("get_slide", {"uid": uid})
        )
        assert moved["start_sec"] == resized_end
        assert moved["end_sec"] == first["end_sec"]
        persisted = next(
            slide
            for slide in json.loads(
                (project_dir / "slides.json").read_text("utf-8")
            )["slides"]
            if slide["uid"] == uid
        )
        assert persisted["start"] == resized_end
        assert persisted["end"] == first["end_sec"]

        _call_and_wait(client, "slide_set_frame", {"uid": uid})
        with_frame = McpClient.result_data(
            client.tool_call("get_slide", {"uid": uid})
        )
        assert with_frame["image_path"]
        assert (project_dir / "slides" / f"slide_{with_frame['index']:03d}.png").is_file()

        _call_and_wait(
            client,
            "slide_clear_image",
            {"uid": uid, "confirm": True},
        )
        cleared = McpClient.result_data(
            client.tool_call("get_slide", {"uid": uid})
        )
        assert cleared["image_path"] == ""

        _call_and_wait(
            client,
            "slide_delete",
            {"uid": uid, "confirm": True},
        )
        project = McpClient.result_data(client.tool_call("get_project"))
        timeline = McpClient.result_data(client.tool_call("get_timeline"))
        document = json.loads((project_dir / "slides.json").read_text("utf-8"))
        project_doc = json.loads((project_dir / "project.json").read_text("utf-8"))
        assert project["slides_count"] == original_count
        assert timeline["tracks"]["slides"]["clip_count"] == original_count
        assert len(document["slides"]) == original_count
        assert len(project_doc["slides"]) == original_count
        assert all(slide["uid"] != uid for slide in document["slides"])
    finally:
        harness.stop()


def test_real_gui_mcp_save_close_open_preserves_project_state(
    repo_dir,
    tmp_path,
    synthetic_video_path,
    synthetic_subtitle_path,
):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"
        project_dir = project_parent / "roundtrip_characterized"
        _call_and_wait(
            client,
            "project_create",
            {"path": str(project_parent), "name": "roundtrip_characterized"},
        )
        _call_and_wait(
            client,
            "video_import",
            {"path": str(synthetic_video_path)},
        )
        _call_and_wait(
            client,
            "subtitle_import",
            {"path": str(synthetic_subtitle_path)},
        )
        _call_and_wait(client, "detect", {"confirm": True}, timeout=120)

        timeline = McpClient.result_data(client.tool_call("get_timeline"))
        first = timeline["tracks"]["slides"]["clips"][0]
        added = _call_and_wait(
            client,
            "slide_add",
            {"ts": (first["start_sec"] + first["end_sec"]) / 2},
        )
        added_uid = added["result"]["uid"]
        _call_and_wait(client, "project_save", {})

        project_before = McpClient.result_data(
            client.tool_call("get_project")
        )
        timeline_before = _normalized_timeline_state(
            McpClient.result_data(client.tool_call("get_timeline"))
        )
        project_json_before = json.loads(
            (project_dir / "project.json").read_text("utf-8")
        )
        slides_json_before = json.loads(
            (project_dir / "slides.json").read_text("utf-8")
        )

        _call_and_wait(
            client,
            "project_close",
            {"confirm": True},
        )
        closed_project = McpClient.result_data(
            client.tool_call("get_project")
        )
        closed_timeline = McpClient.result_data(
            client.tool_call("get_timeline")
        )
        closed_ui = McpClient.result_data(client.tool_call("get_ui_state"))
        assert closed_project["error"] == "no project"
        assert closed_timeline["tracks"] == {}
        for button in (
            "detect",
            "quick_preview",
            "auto",
            "auto_align",
            "export_md",
            "export_pptx",
            "process_notes",
            "save",
        ):
            assert closed_ui["buttons"][button]["enabled"] is False

        _call_and_wait(client, "project_open", {"path": str(project_dir)})
        project_after = McpClient.result_data(client.tool_call("get_project"))
        timeline_after = _normalized_timeline_state(
            McpClient.result_data(client.tool_call("get_timeline"))
        )
        assert project_after == project_before
        assert timeline_after == timeline_before
        assert json.loads(
            (project_dir / "project.json").read_text("utf-8")
        ) == project_json_before
        assert json.loads(
            (project_dir / "slides.json").read_text("utf-8")
        ) == slides_json_before
        assert any(
            slide["uid"] == added_uid for slide in slides_json_before["slides"]
        )
        assert all(
            not slide["image"]
            or (project_dir / slide["image"]).is_file()
            for slide in slides_json_before["slides"]
        )
    finally:
        harness.stop()


def test_real_gui_mcp_markdown_and_pptx_exports_are_structurally_valid(
    repo_dir,
    tmp_path,
    synthetic_video_path,
    synthetic_subtitle_path,
):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"
        project_name = "exports_characterized"
        project_dir = project_parent / project_name
        _call_and_wait(
            client,
            "project_create",
            {"path": str(project_parent), "name": project_name},
        )
        _call_and_wait(
            client,
            "video_import",
            {"path": str(synthetic_video_path)},
        )
        _call_and_wait(
            client,
            "subtitle_import",
            {"path": str(synthetic_subtitle_path)},
        )
        _call_and_wait(client, "detect", {"confirm": True}, timeout=120)

        source = json.loads((project_dir / "slides.json").read_text("utf-8"))
        source_slides = source["slides"]
        assert source_slides
        assert all(slide["image"] for slide in source_slides)

        _call_and_wait(
            client,
            "export_md",
            {"overwrite": True, "confirm": True},
        )
        markdown_path = project_dir / "deck.md"
        markdown = markdown_path.read_text(encoding="utf-8")
        image_refs = re.findall(r"!\[bg\]\(([^)]+)\)", markdown)
        assert "marp: true" in markdown
        assert f'title: "{project_name}"' in markdown
        assert "slides/slides/" not in markdown
        assert len(image_refs) == len(source_slides)
        assert all((project_dir / ref).is_file() for ref in image_refs)
        assert len(
            re.findall(r"^> \d+:\d{2} – \d+:\d{2}$", markdown, re.MULTILINE)
        ) == len(source_slides)

        _call_and_wait(
            client,
            "export_pptx",
            {"overwrite": True, "confirm": True},
        )
        pptx_path = project_dir / "deck.pptx"
        presentation = Presentation(str(pptx_path))
        assert presentation.core_properties.title == project_name
        assert len(presentation.slides) == len(source_slides)
        assert all(
            any(
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                for shape in slide.shapes
            )
            for slide in presentation.slides
        )
        assert all(
            slide.notes_slide.notes_text_frame.text.strip()
            for slide in presentation.slides
        )

        project = McpClient.result_data(client.tool_call("get_project"))
        assert project["pipeline_state"]["md_exported"] is True
        assert project["pipeline_state"]["pptx_exported"] is True
        assert project["artifact_paths"]["deck_md"] == str(markdown_path)
        assert project["artifact_paths"]["deck_pptx"] == str(pptx_path)
    finally:
        harness.stop()


def test_real_gui_mcp_video_and_subtitle_import_gates(
    repo_dir,
    tmp_path,
    synthetic_video_path,
    synthetic_subtitle_path,
):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"
        _call_and_wait(
            client,
            "project_create",
            {"path": str(project_parent), "name": "import_char"},
        )
        _call_and_wait(
            client,
            "video_import",
            {"path": str(synthetic_video_path)},
        )
        project = McpClient.result_data(client.tool_call("get_project"))
        assert project["video"] == str(synthetic_video_path)
        ui = McpClient.result_data(client.tool_call("get_ui_state"))
        assert ui["buttons"]["detect"]["enabled"] is True
        assert ui["buttons"]["quick_preview"]["enabled"] is True

        _call_and_wait(
            client,
            "subtitle_import",
            {"path": str(synthetic_subtitle_path)},
        )
        project = McpClient.result_data(client.tool_call("get_project"))
        timeline = McpClient.result_data(client.tool_call("get_timeline"))
        assert project["subtitle_path"] == str(synthetic_subtitle_path)
        sub_clips = timeline["tracks"].get("subtitles", {}).get("clips", [])
        assert len(sub_clips) == 4
        assert all(clip["text"].strip() for clip in sub_clips)
        assert sub_clips[0]["start_sec"] == 0.0
        assert sub_clips[-1]["end_sec"] == 12.0
    finally:
        harness.stop()


def test_real_gui_mcp_process_notes_enriches_transcript(
    repo_dir,
    tmp_path,
    synthetic_video_path,
    synthetic_subtitle_path,
):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"
        project_dir = project_parent / "notes_char"
        _call_and_wait(
            client,
            "project_create",
            {"path": str(project_parent), "name": "notes_char"},
        )
        _call_and_wait(
            client,
            "video_import",
            {"path": str(synthetic_video_path)},
        )
        _call_and_wait(
            client,
            "subtitle_import",
            {"path": str(synthetic_subtitle_path)},
        )
        _call_and_wait(client, "detect", {"confirm": True}, timeout=120)
        _call_and_wait(
            client,
            "auto_align",
            {"dry_run": False, "max_shift_sec": 3.0, "confirm": True},
        )

        before = json.loads((project_dir / "slides.json").read_text("utf-8"))
        assert not all(slide.get("transcript", "").strip() for slide in before["slides"])

        _call_and_wait(
            client,
            "process_notes",
            {"confirm": True},
            timeout=60,
        )
        after = json.loads((project_dir / "slides.json").read_text("utf-8"))
        assert after["slides"]
        assert any(
            slide.get("transcript", "").strip() for slide in after["slides"]
        )

        project = McpClient.result_data(client.tool_call("get_project"))
        assert project["pipeline_state"]["notes_done"] is True
    finally:
        harness.stop()


def test_real_gui_mcp_full_auto_completes_all_stages(
    repo_dir,
    tmp_path,
    synthetic_video_path,
    synthetic_subtitle_path,
):
    harness = GuiProcessHarness(
        repo=repo_dir,
        run_dir=tmp_path / "run",
        startup_timeout=60.0,
        qt_platform="offscreen",
    )
    try:
        client = harness.start()
        project_parent = tmp_path / "projects"
        project_dir = project_parent / "auto_char"
        _call_and_wait(
            client,
            "project_create",
            {"path": str(project_parent), "name": "auto_char"},
        )
        _call_and_wait(
            client,
            "video_import",
            {"path": str(synthetic_video_path)},
        )
        _call_and_wait(
            client,
            "subtitle_import",
            {"path": str(synthetic_subtitle_path)},
        )

        _call_and_wait(client, "auto", {"confirm": True}, timeout=180)

        project = McpClient.result_data(client.tool_call("get_project"))
        state = project["pipeline_state"]
        assert state["detect_done"] is True
        assert state["align_done"] is True
        assert state["notes_done"] is True
        assert state["md_exported"] is True
        assert state["pptx_exported"] is True
        assert state["auto_done"] is True
        assert project["slides_count"] > 0
        assert (project_dir / "slides.json").is_file()
        assert (project_dir / "deck.md").is_file()
        assert (project_dir / "deck.pptx").is_file()
        assert (project_dir / "alignment_report.json").is_file()
        assert (project_dir / "slides").is_dir()
    finally:
        harness.stop()
