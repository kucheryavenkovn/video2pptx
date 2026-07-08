# FILE: tests/test_pptx_export.py
# VERSION: 0.1.0
# START_MODULE_CONTRACT
#   PURPOSE: Unit tests for pptx_export module
#   SCOPE: Verify export_to_pptx creates valid .pptx with correct slide count,
#          notes with formatted timestamps/cues, and embedded images
#   DEPENDS: pytest, pptx_export, models
#   LINKS: M-PPTX-EXPORT, V-M-PPTX-EXPORT
#   ROLE: TEST
#   MAP_MODE: EXPORTS
# END_MODULE_CONTRACT

from __future__ import annotations

import json
from pathlib import Path

import pytest

from video_slide_md.models import SlidesDocument, SlideSegment, SubtitleCue, VideoInfo
from video_slide_md.pptx_export import export_to_pptx, _format_slide_notes, _group_cues


class TestGroupCues:
    def test_empty_cues(self):
        assert _group_cues([]) == []

    def test_single_cue(self):
        cues = [SubtitleCue(start=0.0, end=5.0, text="Hello")]
        result = _group_cues(cues)
        assert len(result) == 1
        assert result[0] == cues

    def test_continuous_cues_grouped(self):
        cues = [
            SubtitleCue(start=0.0, end=2.0, text="First"),
            SubtitleCue(start=2.5, end=5.0, text="Second"),
            SubtitleCue(start=5.5, end=8.0, text="Third"),
        ]
        result = _group_cues(cues, gap_threshold=2.0)
        assert len(result) == 1

    def test_gapped_cues_separate(self):
        cues = [
            SubtitleCue(start=0.0, end=2.0, text="First"),
            SubtitleCue(start=5.0, end=8.0, text="Second"),
        ]
        result = _group_cues(cues, gap_threshold=2.0)
        assert len(result) == 2


class TestFormatSlideNotes:
    def test_with_cues(self):
        seg = SlideSegment(
            index=1,
            start=0.0,
            end=10.0,
            duration=10.0,
            representative_timestamp=5.0,
            subtitle_cues=[
                SubtitleCue(start=0.0, end=3.0, text="Hello there"),
                SubtitleCue(start=3.5, end=6.0, text="How are you"),
            ],
        )
        result = _format_slide_notes(seg)
        assert "[ 0:00 – 0:10 ]" in result
        assert "Hello there" in result
        assert "How are you" in result

    def test_no_cues_fallback_transcript(self):
        seg = SlideSegment(
            index=1,
            start=0.0,
            end=5.0,
            duration=5.0,
            representative_timestamp=2.5,
            transcript="Just some transcript text.",
        )
        result = _format_slide_notes(seg)
        assert "Just some transcript text." in result

    def test_no_cues_no_transcript(self):
        seg = SlideSegment(index=1, start=0.0, end=5.0, duration=5.0, representative_timestamp=2.5)
        result = _format_slide_notes(seg)
        assert result == "[ 0:00 \u2013 0:05 ]"


class TestExportToPptx:
    def test_minimal_document(self, tmp_path: Path):
        doc = SlidesDocument(
            video=VideoInfo(path="test.mp4", duration=10.0, width=1920, height=1080, fps=30.0),
            slides=[],
        )
        out = tmp_path / "test.pptx"
        result = export_to_pptx(doc, out)
        assert result == out
        assert out.is_file()
        assert out.stat().st_size > 0

    def test_slides_with_notes(self, tmp_path: Path):
        from pptx import Presentation as PptxPresentation

        doc = SlidesDocument(
            video=VideoInfo(path="test.mp4", duration=30.0, width=1920, height=1080, fps=30.0),
            slides=[
                SlideSegment(
                    index=1,
                    start=0.0,
                    end=10.0,
                    duration=10.0,
                    representative_timestamp=5.0,
                    transcript="First slide text.",
                    subtitle_cues=[
                        SubtitleCue(start=0.0, end=5.0, text="First cue"),
                        SubtitleCue(start=5.5, end=10.0, text="Second cue"),
                    ],
                ),
                SlideSegment(
                    index=2,
                    start=10.0,
                    end=20.0,
                    duration=10.0,
                    representative_timestamp=15.0,
                    transcript="Second slide text.",
                    subtitle_cues=[
                        SubtitleCue(start=10.0, end=15.0, text="Third cue"),
                    ],
                ),
            ],
        )
        out = tmp_path / "test.pptx"
        export_to_pptx(doc, out, slides_dir=tmp_path)

        prs = PptxPresentation(str(out))
        assert len(prs.slides) == 2

        # Check notes on first slide
        notes0 = prs.slides[0].notes_slide.notes_text_frame.text
        assert "[ 0:00 – 0:10 ]" in notes0
        assert "First cue" in notes0
        assert "Second cue" in notes0

        notes1 = prs.slides[1].notes_slide.notes_text_frame.text
        assert "[ 0:10 – 0:20 ]" in notes1
        assert "Third cue" in notes1

    def test_no_notes_does_not_crash(self, tmp_path: Path):
        doc = SlidesDocument(
            video=VideoInfo(path="test.mp4", duration=5.0, width=640, height=480, fps=30.0),
            slides=[
                SlideSegment(index=1, start=0.0, end=5.0, duration=5.0, representative_timestamp=2.5),
            ],
        )
        out = tmp_path / "test.pptx"
        result = export_to_pptx(doc, out)
        assert out.is_file()
